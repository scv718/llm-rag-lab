# app/services/zip_ingest.py
from __future__ import annotations

from dataclasses import dataclass
import csv
import io
from pathlib import Path
from typing import Iterable, List, Tuple
import os
import re
import zipfile
import hashlib


TEXT_EXTS = {
    ".java", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".md", ".txt",
    ".yml", ".yaml",
    ".json", ".xml",
    ".properties", ".gradle", ".sql",
    ".sh", ".bat", ".ps1",
    ".ini", ".cfg",
}

ARCHIVE_DOC_EXTS = {
    ".txt", ".md", ".csv", ".pdf", ".docx", ".xlsx", ".pptx",
}

# 업로드 zip 내부에서 흔히 섞이는 것들 제외
SKIP_DIR_NAMES = {".git", ".idea", ".vscode", "__pycache__", "node_modules", "dist", "build", "target", "venv"}
SKIP_FILE_NAMES = {".env", ".env.local", ".env.dev", ".env.prod"}  # 키 유출 방지


@dataclass(frozen=True)
class CodeChunk:
    chunk_id: str
    repo_id: str
    path: str
    lang: str
    start_line: int
    end_line: int
    text: str


def sha256_bytes(b: bytes) -> str:
    return hashlib.sha256(b).hexdigest()


def _is_within_dir(base: Path, target: Path) -> bool:
    # Zip Slip 방어: base 밖으로 나가면 False
    try:
        target.resolve().relative_to(base.resolve())
        return True
    except Exception:
        return False


def safe_extract_zip(
    *,
    zip_path: Path,
    extract_dir: Path,
    max_files: int = 5000,
    max_total_uncompressed_bytes: int = 200 * 1024 * 1024,  # 200MB
) -> Tuple[int, int]:
    """
    안전한 압축 해제:
    - Zip Slip 방어
    - 파일 개수 제한
    - 해제 총량 제한
    """
    extract_dir.mkdir(parents=True, exist_ok=True)

    total = 0
    count = 0

    with zipfile.ZipFile(zip_path, "r") as zf:
        infos = zf.infolist()
        if len(infos) > max_files:
            raise ValueError(f"ZIP contains too many entries: {len(infos)} > {max_files}")

        for info in infos:
            # 디렉토리 엔트리 스킵
            if info.is_dir():
                continue

            # 총 해제 바이트 제한(압축폭탄 방어의 1차)
            total += int(info.file_size or 0)
            if total > max_total_uncompressed_bytes:
                raise ValueError("ZIP uncompressed size limit exceeded")

            # 경로 정규화 + Zip Slip 방어
            rel = Path(info.filename)
            # zip 내부 경로는 POSIX 스타일이 많음. Path가 알아서 처리하긴 하지만, 절대경로/드라이브 제거 방어
            rel = Path(*[p for p in rel.parts if p not in ("", ".", "..")])
            out_path = extract_dir / rel

            if not _is_within_dir(extract_dir, out_path):
                raise ValueError("Blocked Zip Slip attempt")

            out_path.parent.mkdir(parents=True, exist_ok=True)
            with zf.open(info, "r") as src, open(out_path, "wb") as dst:
                dst.write(src.read())

            count += 1

    return count, total


def _guess_lang(path: str) -> str:
    p = Path(path)
    if p.name == "Dockerfile":
        return "dockerfile"
    ext = p.suffix.lower()
    if ext == ".csv":
        return "csv"
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".py":
        return "python"
    if ext == ".java":
        return "java"
    if ext in (".js", ".jsx"):
        return "javascript"
    if ext in (".ts", ".tsx"):
        return "typescript"
    if ext in (".yml", ".yaml"):
        return "yaml"
    if ext == ".json":
        return "json"
    if ext == ".xml":
        return "xml"
    if ext == ".md":
        return "markdown"
    if ext == ".sql":
        return "sql"
    return ext.lstrip(".") or "text"


def iter_source_files(root: Path) -> Iterable[Path]:
    for p in root.rglob("*"):
        if p.is_dir():
            continue

        # 스킵 디렉토리 포함이면 제외
        if any(part in SKIP_DIR_NAMES for part in p.parts):
            continue

        if p.name in SKIP_FILE_NAMES:
            continue

        if p.name == "Dockerfile":
            yield p
            continue

        ext = p.suffix.lower()
        if ext in TEXT_EXTS or ext in ARCHIVE_DOC_EXTS:
            yield p


def read_text_best_effort(path: Path, max_bytes: int = 2 * 1024 * 1024) -> str:
    """
    텍스트 파일 best-effort 읽기
    - 너무 큰 파일은 상한으로 컷
    - 인코딩은 utf-8 우선, 실패 시 cp949, latin-1 순으로 시도
    """
    raw = path.read_bytes()
    return decode_text_best_effort(raw, max_bytes=max_bytes)


def decode_text_best_effort(raw: bytes, max_bytes: int = 2 * 1024 * 1024) -> str:
    """
    바이트열을 텍스트로 best-effort 디코딩한다.
    """
    if len(raw) > max_bytes:
        raw = raw[:max_bytes]

    for enc in ("utf-8", "utf-8-sig", "cp949", "latin-1"):
        try:
            return raw.decode(enc)
        except Exception:
            continue
    # 최후: 에러 무시
    return raw.decode("utf-8", errors="ignore")


def extract_text_from_path(path: Path, max_bytes: int = 2 * 1024 * 1024) -> str:
    ext = path.suffix.lower()

    if ext in TEXT_EXTS:
        return read_text_best_effort(path, max_bytes=max_bytes)
    if ext == ".csv":
        return extract_csv_text(path, max_bytes=max_bytes)
    if ext == ".pdf":
        return extract_pdf_text(path, max_bytes=max_bytes)
    if ext == ".docx":
        return extract_docx_text(path)
    if ext == ".xlsx":
        return extract_xlsx_text(path)
    if ext == ".pptx":
        return extract_pptx_text(path)
    return ""


def extract_csv_text(path: Path, max_bytes: int = 2 * 1024 * 1024) -> str:
    text = read_text_best_effort(path, max_bytes=max_bytes)
    rows = list(csv.reader(io.StringIO(text)))
    if not rows:
        return ""

    header = rows[0]
    lines = []
    for row_idx, row in enumerate(rows[1:], start=2):
        pairs = []
        for col_idx, value in enumerate(row):
            key = header[col_idx].strip() if col_idx < len(header) and header[col_idx].strip() else f"col_{col_idx + 1}"
            pairs.append(f"{key}: {value}")
        if pairs:
            lines.append(f"row {row_idx} | " + " | ".join(pairs))
    return "\n".join(lines) if lines else text


def extract_pdf_text(path: Path, max_bytes: int = 10 * 1024 * 1024) -> str:
    try:
        import PyPDF2
    except ImportError:
        return ""

    raw = path.read_bytes()
    if len(raw) > max_bytes:
        raw = raw[:max_bytes]

    reader = PyPDF2.PdfReader(io.BytesIO(raw))
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    return "\n\n".join([page for page in pages if page])


def extract_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""

    document = Document(str(path))
    blocks = []
    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            blocks.append(text)
    for table in document.tables:
        for row_idx, row in enumerate(table.rows, start=1):
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                blocks.append(f"table row {row_idx} | " + " | ".join(cells))
    return "\n\n".join(blocks)


def extract_xlsx_text(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""

    workbook = load_workbook(path, data_only=True, read_only=True)
    lines = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        header = None
        for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = ["" if value is None else str(value).strip() for value in row]
            if not any(values):
                continue
            if header is None:
                header = values
                continue
            pairs = []
            for col_idx, value in enumerate(values):
                if not value:
                    continue
                key = header[col_idx].strip() if col_idx < len(header) and header[col_idx].strip() else f"col_{col_idx + 1}"
                pairs.append(f"{key}: {value}")
            if pairs:
                lines.append(f"sheet {sheet_name} row {row_idx} | " + " | ".join(pairs))
    workbook.close()
    return "\n".join(lines)


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    presentation = Presentation(str(path))
    lines = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                texts.append(text)
        if texts:
            lines.append(f"slide {slide_idx} | " + " | ".join(texts))
    return "\n".join(lines)


def chunk_code_by_lines(
    *,
    repo_id: str,
    rel_path: str,
    text: str,
    lines_per_chunk: int = 250,
    overlap: int = 40,
) -> List[CodeChunk]:
    if lines_per_chunk <= 0:
        raise ValueError("lines_per_chunk must be > 0")
    if overlap < 0 or overlap >= lines_per_chunk:
        raise ValueError("overlap must be >=0 and < lines_per_chunk")

    lines = text.splitlines()
    if not lines:
        return []

    lang = _guess_lang(rel_path)
    boundaries = detect_code_boundaries(lines, lang)

    chunks: List[CodeChunk] = []
    step = lines_per_chunk - overlap
    start = 0
    idx = 0

    while start < len(lines):
        target_end = min(start + lines_per_chunk, len(lines))
        end = choose_chunk_end(lines, start, target_end, boundaries)
        piece_lines = lines[start:end]
        piece = "\n".join(piece_lines).strip()

        if piece:
            chunk_id = f"{repo_id}:{rel_path}:L{start+1}-L{end}:{idx:03d}"
            chunks.append(
                CodeChunk(
                    chunk_id=chunk_id,
                    repo_id=repo_id,
                    path=rel_path,
                    lang=lang,
                    start_line=start + 1,
                    end_line=end,
                    text=piece,
                )
            )
            idx += 1
            if end >= len(lines):
                break

            next_start = max(start + step, end - overlap)
            if next_start <= start:
                next_start = end
            start = next_start
            continue

        start += 1

    return chunks


def detect_code_boundaries(lines: List[str], lang: str) -> set[int]:
    boundaries = {0, len(lines)}

    for idx, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            boundaries.add(idx)
            continue

        if looks_like_definition_boundary(stripped, lang):
            boundaries.add(idx)

    return boundaries


def looks_like_definition_boundary(stripped: str, lang: str) -> bool:
    if lang in {"python"}:
        return bool(re.match(r"^(class|def|async def)\s+", stripped))

    if lang in {"java", "javascript", "typescript", "go", "cs", "cpp", "c", "php", "kt"}:
        if re.match(r"^(public|private|protected|internal|static|final|abstract|sealed|\@)", stripped):
            return True
        if re.match(r"^(class|interface|enum|record|object|data class|fun|func)\b", stripped):
            return True
        if re.match(r"^[A-Za-z0-9_<>\[\], ?]+\s+[A-Za-z_][A-Za-z0-9_]*\s*\(", stripped):
            return True

    if lang in {"yaml", "json", "xml", "sql", "markdown", "properties"}:
        return bool(re.match(r"^([A-Za-z0-9_.\-]+:|<[^/!?][^>]*>|(select|insert|update|delete|create|alter)\b|#+\s)", stripped, re.IGNORECASE))

    return False


def choose_chunk_end(lines: List[str], start: int, target_end: int, boundaries: set[int], lookaround: int = 40) -> int:
    if target_end >= len(lines):
        return len(lines)

    candidate_end = target_end
    best_score = boundary_score(lines, candidate_end, boundaries)

    lower = max(start + 1, target_end - lookaround)
    upper = min(len(lines), target_end + lookaround)
    for end in range(lower, upper + 1):
        score = boundary_score(lines, end, boundaries)
        distance_penalty = abs(end - target_end) * 0.03
        final_score = score - distance_penalty
        if final_score > best_score:
            best_score = final_score
            candidate_end = end

    return max(start + 1, candidate_end)


def boundary_score(lines: List[str], end: int, boundaries: set[int]) -> float:
    score = 0.0
    if end in boundaries:
        score += 2.0

    prev_line = lines[end - 1].strip() if 0 < end <= len(lines) else ""
    next_line = lines[end].strip() if end < len(lines) else ""

    if not prev_line or not next_line:
        score += 1.0
    if prev_line.endswith(("}", "];", ");", ")", ":")):
        score += 0.4
    if next_line.startswith(("class ", "def ", "async def ", "public ", "private ", "@", "function ")):
        score += 0.8
    return score
